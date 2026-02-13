#!/usr/bin/env python3
"""
Generate the PepsiCo RLM AI Observability Investigation Layer deck (v2).

Design system: "Evidence Ledger" style
- Warm paper background, deep ink navy text, monospace evidence chips
- 60/40 grid, card system, budget strips, run-record stamps

Usage:
    pip install python-pptx
    python generate_deck_v2.py
    # -> outputs ppevigil_rlm_pepsico_deck_v2.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ---------------------------------------------------------------------------
# Design tokens
# ---------------------------------------------------------------------------
BG_WARM_PAPER = RGBColor(0xFB, 0xF7, 0xF0)
INK_NAVY = RGBColor(0x0B, 0x1F, 0x33)
SLATE = RGBColor(0x54, 0x65, 0x7A)
BORDER_TAN = RGBColor(0xD7, 0xCB, 0xBE)
ACCENT_RED = RGBColor(0xC6, 0x28, 0x28)
TEAL = RGBColor(0x1B, 0x7F, 0x79)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
AMBER = RGBColor(0xC4, 0x7F, 0x17)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF2, 0xEE, 0xE6)
CARD_BG = RGBColor(0xFF, 0xFD, 0xFA)
MONO_BG = RGBColor(0xF5, 0xF0, 0xE8)

FONT_TITLE = "Segoe UI Semibold"
FONT_BODY = "Segoe UI"
FONT_MONO = "Consolas"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.5)

TOTAL_SLIDES = 12

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def set_slide_bg(slide, color=BG_WARM_PAPER):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_name=FONT_BODY,
                font_size=Pt(14), font_color=INK_NAVY, bold=False, italic=False,
                alignment=PP_ALIGN.LEFT, word_wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = font_size
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.italic = italic
    p.alignment = alignment
    return txBox


def add_rich_textbox(slide, left, top, width, height, runs,
                     alignment=PP_ALIGN.LEFT, line_spacing=1.15):
    """runs = list of dicts: {text, font_name, font_size, font_color, bold, italic}"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    if line_spacing:
        p.line_spacing = line_spacing
    for i, r in enumerate(runs):
        if i == 0:
            run = p.runs[0] if p.runs else p.add_run()
        else:
            run = p.add_run()
        run.text = r.get("text", "")
        run.font.name = r.get("font_name", FONT_BODY)
        run.font.size = r.get("font_size", Pt(14))
        run.font.color.rgb = r.get("font_color", INK_NAVY)
        run.font.bold = r.get("bold", False)
        run.font.italic = r.get("italic", False)
    return txBox


def add_bullet_slide_text(slide, left, top, width, height, lines,
                           font_name=FONT_BODY, font_size=Pt(13),
                           font_color=INK_NAVY, level_indent=Inches(0.3)):
    """lines = list of (level:int, text:str) or (level, text, overrides_dict)."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        level = item[0]
        text = item[1]
        overrides = item[2] if len(item) > 2 else {}
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.level = level
        p.font.name = overrides.get("font_name", font_name)
        p.font.size = overrides.get("font_size", font_size)
        p.font.color.rgb = overrides.get("font_color", font_color)
        p.font.bold = overrides.get("bold", False)
        p.space_after = Pt(4)
        if level > 0:
            p.font.size = Pt(font_size.pt - 1)
    return txBox


def add_card(slide, left, top, width, height, fill_color=CARD_BG,
             border_color=BORDER_TAN, border_width=Pt(1), corner_radius=Inches(0.08)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = border_color
    shape.line.width = border_width
    # Adjust corner radius via adjustments
    if shape.adjustments:
        try:
            shape.adjustments[0] = 0.04  # ~4% corner radius
        except Exception:
            pass
    return shape


def add_chip(slide, left, top, text, bg_color=MONO_BG, text_color=INK_NAVY,
             font_size=Pt(9), width=None):
    """Small rounded pill with monospace text."""
    if width is None:
        width = Inches(max(1.2, len(text) * 0.085))
    height = Inches(0.28)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = BORDER_TAN
    shape.line.width = Pt(0.5)
    tf = shape.text_frame
    tf.word_wrap = False
    tf.margin_left = Pt(6)
    tf.margin_right = Pt(6)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = FONT_MONO
    p.font.size = font_size
    p.font.color.rgb = text_color
    p.alignment = PP_ALIGN.CENTER
    return shape


def add_stamp(slide, left, top, text="RUN RECORDED", color=GREEN):
    width = Inches(1.3)
    height = Inches(0.26)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = color
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.word_wrap = False
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = FONT_MONO
    p.font.size = Pt(8)
    p.font.color.rgb = color
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    return shape


def add_budget_strip(slide, left, top, text="depth 1/2 | iter 2/8 | tools 3/16 | cost $0.01 | wall 12s",
                     width=Inches(4.6)):
    height = Inches(0.3)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = MONO_BG
    shape.line.color.rgb = BORDER_TAN
    shape.line.width = Pt(0.75)
    tf = shape.text_frame
    tf.word_wrap = False
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = FONT_MONO
    p.font.size = Pt(8)
    p.font.color.rgb = SLATE
    p.alignment = PP_ALIGN.LEFT
    return shape


def add_ledger_header(slide, text="AI Observability Investigation Layer"):
    """Small top-left label with thin vertical red rule."""
    # Red vertical rule
    left = MARGIN
    top = Inches(0.3)
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, Pt(4), Inches(0.3)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_RED
    line.line.fill.background()
    # Label
    add_textbox(slide, left + Pt(12), top, Inches(3.5), Inches(0.3),
                text, font_name=FONT_BODY, font_size=Pt(9),
                font_color=SLATE)


def add_footer(slide, slide_num):
    text = f"Slide {slide_num}/{TOTAL_SLIDES}"
    add_textbox(slide, Inches(11.5), Inches(7.05), Inches(1.5), Inches(0.3),
                text, font_name=FONT_BODY, font_size=Pt(8),
                font_color=SLATE, alignment=PP_ALIGN.RIGHT)


def add_slide_title(slide, title, subtitle=None, left=MARGIN, top=Inches(0.75)):
    add_textbox(slide, left, top, Inches(7.5), Inches(0.6),
                title, font_name=FONT_TITLE, font_size=Pt(26),
                font_color=INK_NAVY, bold=True)
    if subtitle:
        add_textbox(slide, left, top + Inches(0.6), Inches(7.5), Inches(0.4),
                    subtitle, font_name=FONT_BODY, font_size=Pt(14),
                    font_color=SLATE)


def add_arrow(slide, start_left, start_top, end_left, end_top, color=BORDER_TAN, width=Pt(1.5)):
    """Add a connector arrow."""
    connector = slide.shapes.add_connector(
        1,  # straight connector
        start_left, start_top, end_left, end_top
    )
    connector.line.color.rgb = color
    connector.line.width = width
    return connector


def add_step_card(slide, left, top, num, label, detail, width=Inches(4.5)):
    """A numbered step card for the runtime ledger."""
    height = Inches(0.48)
    card = add_card(slide, left, top, width, height)
    # Number gutter
    add_textbox(slide, left + Pt(6), top + Pt(2), Inches(0.35), Inches(0.35),
                f"{num:02d}", font_name=FONT_MONO, font_size=Pt(9),
                font_color=SLATE, bold=True)
    # Label
    add_textbox(slide, left + Inches(0.4), top + Pt(2), Inches(1.3), Inches(0.35),
                label, font_name=FONT_TITLE, font_size=Pt(11),
                font_color=INK_NAVY, bold=True)
    # Detail
    add_textbox(slide, left + Inches(1.7), top + Pt(3), width - Inches(1.8), Inches(0.35),
                detail, font_name=FONT_MONO, font_size=Pt(8),
                font_color=SLATE)
    return card


def add_speaker_notes(slide, text):
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text


def add_simple_table(slide, left, top, width, height, rows, col_widths=None,
                     header_bg=MONO_BG, header_color=INK_NAVY):
    """rows = list of lists of strings. First row is header."""
    n_rows = len(rows)
    n_cols = len(rows[0]) if rows else 0
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for r_idx, row in enumerate(rows):
        for c_idx, cell_text in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = cell_text
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.name = FONT_MONO if r_idx > 0 else FONT_TITLE
                paragraph.font.size = Pt(9)
                paragraph.font.color.rgb = header_color if r_idx == 0 else INK_NAVY
                paragraph.font.bold = (r_idx == 0)
            # Header row styling
            if r_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_bg
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = CARD_BG if r_idx % 2 == 1 else WHITE

            # Borders
            from pptx.oxml.ns import qn
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            for edge in ['a:lnL', 'a:lnR', 'a:lnT', 'a:lnB']:
                ln = tcPr.find(qn(edge))
                if ln is None:
                    from lxml import etree
                    ln = etree.SubElement(tcPr, qn(edge))
                ln.set('w', '6350')  # 0.5pt
                sf = ln.find(qn('a:solidFill'))
                if sf is None:
                    from lxml import etree
                    sf = etree.SubElement(ln, qn('a:solidFill'))
                srgb = sf.find(qn('a:srgbClr'))
                if srgb is None:
                    from lxml import etree
                    srgb = etree.SubElement(sf, qn('a:srgbClr'))
                srgb.set('val', 'D7CBBE')

    return table_shape


# ---------------------------------------------------------------------------
# Faint ledger lines on background
# ---------------------------------------------------------------------------
def add_ledger_lines(slide, count=12):
    """Add faint horizontal lines across the slide for the ledger effect."""
    spacing = SLIDE_H // (count + 1)
    for i in range(1, count + 1):
        y = spacing * i
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), y, SLIDE_W, Pt(0.5)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = BORDER_TAN
        # Make it very faint via transparency on the shape
        from pptx.oxml.ns import qn
        spPr = line._element.find(qn('p:spPr'))
        if spPr is None:
            spPr = line._element.find(qn('p:spPr'))
        solidFill = spPr.find(qn('a:solidFill'))
        if solidFill is not None:
            srgb = solidFill.find(qn('a:srgbClr'))
            if srgb is not None:
                from lxml import etree
                alpha = etree.SubElement(srgb, qn('a:alpha'))
                alpha.set('val', '12000')  # 12% opacity
        line.line.fill.background()


# ===========================================================================
# SLIDE BUILDERS
# ===========================================================================

def build_slide_1(prs):
    """Title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide)
    add_ledger_lines(slide)

    # Red vertical rule
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, MARGIN, Inches(1.5), Pt(6), Inches(2.8)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_RED
    line.line.fill.background()

    # Title
    add_textbox(slide, MARGIN + Inches(0.25), Inches(1.6), Inches(7.2), Inches(0.9),
                "AI Observability\nInvestigation Layer",
                font_name=FONT_TITLE, font_size=Pt(36), font_color=INK_NAVY, bold=True)

    # Subtitle: "(RLM-powered)"
    add_textbox(slide, MARGIN + Inches(0.25), Inches(2.8), Inches(7), Inches(0.4),
                "(RLM-powered)",
                font_name=FONT_TITLE, font_size=Pt(20), font_color=TEAL, bold=True)

    # Promise line
    add_textbox(slide, MARGIN + Inches(0.25), Inches(3.35), Inches(7), Inches(0.5),
                "From long-context bottlenecks to low-cost\nautomated RCA and policy compliance",
                font_name=FONT_BODY, font_size=Pt(16), font_color=SLATE)

    # Micro-diagram: Traces -> RLM Runtime -> Findings
    chip_y = Inches(4.3)
    chip1 = add_chip(slide, MARGIN + Inches(0.3), chip_y, "Traces", width=Inches(1.0))
    # Arrow
    add_textbox(slide, MARGIN + Inches(1.4), chip_y - Pt(1), Inches(0.5), Inches(0.28),
                "\u2192", font_name=FONT_MONO, font_size=Pt(14), font_color=SLATE,
                alignment=PP_ALIGN.CENTER)
    chip2 = add_chip(slide, MARGIN + Inches(1.9), chip_y, "RLM Runtime", width=Inches(1.3),
                     bg_color=RGBColor(0xE8, 0xF5, 0xF3), text_color=TEAL)
    add_textbox(slide, MARGIN + Inches(3.3), chip_y - Pt(1), Inches(0.5), Inches(0.28),
                "\u2192", font_name=FONT_MONO, font_size=Pt(14), font_color=SLATE,
                alignment=PP_ALIGN.CENTER)
    chip3 = add_chip(slide, MARGIN + Inches(3.8), chip_y, "Findings", width=Inches(1.0))

    # Right side: faint watermark card
    wm_card = add_card(slide, Inches(8.5), Inches(1.8), Inches(4.0), Inches(3.5),
                       fill_color=RGBColor(0xF8, 0xF4, 0xED),
                       border_color=RGBColor(0xE8, 0xE0, 0xD5))
    # Faint trace tree inside watermark
    for i, label in enumerate(["root_span", "  tool_call", "  retrieval", "    chunk_eval", "  synthesis"]):
        add_textbox(slide, Inches(8.8), Inches(2.1) + Inches(i * 0.5),
                    Inches(3.2), Inches(0.3),
                    label, font_name=FONT_MONO, font_size=Pt(10),
                    font_color=RGBColor(0xD0, 0xC8, 0xBC))

    add_footer(slide, 1)
    add_speaker_notes(slide,
        "This is a recursive, evidence-driven investigation layer on top of our AI "
        "observability stack. It turns messy trace evidence into structured, auditable "
        "RCA and compliance findings.")


def build_slide_2(prs):
    """AI Observability: The Long-Context Problem."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_ledger_lines(slide)
    add_ledger_header(slide)
    add_slide_title(slide, "AI Observability: The Long-Context Problem")

    # Left bullets
    bullets = [
        (0, "AI incidents are not a single log line; evidence is distributed across:"),
        (1, "trace trees (parent/child spans)"),
        (1, "tool I/O artifacts and errors"),
        (1, "retrieval chunks and citations"),
        (1, "configs/snapshots and diffs"),
        (1, "policy controls and required evidence lists"),
        (0, ""),
        (0, '"Long context" fails in practice:'),
        (1, 'the important clue is often buried ("lost in the middle")'),
        (1, "stuffing everything into one prompt increases cost + confusion"),
        (1, "repeated manual triage does not scale"),
        (0, ""),
        (0, "We need automation that is evidence-grounded, budgeted, and replayable.",
         {"bold": True}),
    ]
    add_bullet_slide_text(slide, MARGIN, Inches(1.5), Inches(6.5), Inches(5.0), bullets)

    # Right: Evidence Ledger card
    card_left = Inches(7.5)
    card_top = Inches(1.3)
    card = add_card(slide, card_left, card_top, Inches(5.3), Inches(4.5))

    # Ledger header
    add_textbox(slide, card_left + Pt(12), card_top + Pt(8), Inches(4), Inches(0.3),
                "Evidence Ledger", font_name=FONT_TITLE, font_size=Pt(12),
                font_color=INK_NAVY, bold=True)

    # Table inside card
    ledger_rows = [
        ["Signal", "Contains", "Clue Location"],
        ["trace spans", "parent/child graph", "structure"],
        ["tool I/O", "function inputs/outputs", "errors"],
        ["retrieval chunks", "RAG context", "\u2190 CRITICAL"],
        ["config diffs", "prompt/model changes", "drift"],
        ["controls", "policy requirements", "gaps"],
        ["messages", "user/assistant turns", "intent"],
        ["timeouts", "latency signals", "bottleneck"],
        ["schema errors", "output validation", "contract"],
    ]
    add_simple_table(slide, card_left + Pt(8), card_top + Inches(0.4),
                     Inches(5.0), Inches(3.3), ledger_rows,
                     col_widths=[Inches(1.4), Inches(2.0), Inches(1.6)])

    # "Lost in the middle" callout
    add_textbox(slide, card_left + Inches(3.7), card_top + Inches(1.95),
                Inches(1.5), Inches(0.25),
                "\u2190 lost in the middle",
                font_name=FONT_MONO, font_size=Pt(8), font_color=ACCENT_RED, bold=True)

    # Bottom comparison cards
    # One-shot prompt card
    os_left = Inches(7.7)
    os_top = Inches(6.0)
    add_card(slide, os_left, os_top, Inches(2.3), Inches(0.9),
             fill_color=RGBColor(0xFD, 0xF0, 0xF0), border_color=ACCENT_RED)
    add_textbox(slide, os_left + Pt(8), os_top + Pt(4), Inches(2.1), Inches(0.25),
                "One-shot prompt", font_name=FONT_TITLE, font_size=Pt(9),
                font_color=ACCENT_RED, bold=True)
    add_textbox(slide, os_left + Pt(8), os_top + Pt(22), Inches(2.1), Inches(0.3),
                "\u26a0 tokens \u2191, clarity \u2193",
                font_name=FONT_MONO, font_size=Pt(8), font_color=ACCENT_RED)

    # Bounded investigation card
    bi_left = Inches(10.3)
    add_card(slide, bi_left, os_top, Inches(2.3), Inches(0.9),
             fill_color=RGBColor(0xEE, 0xF8, 0xF0), border_color=GREEN)
    add_textbox(slide, bi_left + Pt(8), os_top + Pt(4), Inches(2.1), Inches(0.25),
                "Bounded investigation", font_name=FONT_TITLE, font_size=Pt(9),
                font_color=GREEN, bold=True)
    add_textbox(slide, bi_left + Pt(8), os_top + Pt(22), Inches(2.1), Inches(0.3),
                "plan \u2192 inspect \u2192 submit",
                font_name=FONT_MONO, font_size=Pt(8), font_color=GREEN)

    # Evidence chips
    chip_y = Inches(5.8)
    for i, label in enumerate(["trace_id: abc-123", "span_id: def-456", "artifact_id: ghi-789"]):
        add_chip(slide, card_left + Inches(i * 1.7) + Pt(8), chip_y, label, font_size=Pt(7))

    add_footer(slide, 2)
    add_speaker_notes(slide,
        'Emphasize that "more tokens" is not the solution: it increases spend and '
        'still misses structured investigation steps.\n'
        'The enterprise pain is "repetitive investigations" and "inconsistent decisions."')


def build_slide_3(prs):
    """What An RLM Is (Practical Definition)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_ledger_lines(slide)
    add_ledger_header(slide)
    add_slide_title(slide, "What An RLM Is (Practical Definition)")

    # Left bullets
    bullets = [
        (0, "An RLM is not just a model call; it's model + runtime:", {"bold": True}),
        (1, "plan next step"),
        (1, "gather evidence via tools"),
        (1, "optionally delegate sub-investigations"),
        (1, "synthesize, then finalize into a schema"),
        (0, ""),
        (0, "Key property: recursive decomposition instead of monolithic prompting.",
         {"bold": True}),
        (0, ""),
        (0, "Guardrails are part of the design:"),
        (1, "budgets (iterations, depth, tools, tokens, cost, wall-time)"),
        (1, "sandbox / tool allowlist"),
        (1, "schema validation + run artifacts"),
    ]
    add_bullet_slide_text(slide, MARGIN, Inches(1.5), Inches(6.5), Inches(5.0), bullets)

    # Right: Runtime ledger (stacked step cards)
    right_left = Inches(7.5)
    steps = [
        (1, "PLAN", 'objective: "label trace failure mode"'),
        (2, "TOOL_CALL", "tool: list_spans(trace_id=...)"),
        (3, "DELEGATE_SUBCALL", 'subcall: hypothesis="retrieval_failure"'),
        (4, "SYNTHESIZE", "patch: add evidence_refs[2]"),
        (5, "FINALIZE", "submit: RCAReport"),
    ]
    for i, (num, label, detail) in enumerate(steps):
        add_step_card(slide, right_left, Inches(1.3) + Inches(i * 0.58),
                      num, label, detail)

    # Budget strip
    add_budget_strip(slide, right_left, Inches(4.3))

    # Return value card
    rv_top = Inches(4.75)
    add_card(slide, right_left, rv_top, Inches(4.5), Inches(1.2),
             fill_color=MONO_BG)
    add_textbox(slide, right_left + Pt(10), rv_top + Pt(6),
                Inches(3.5), Inches(0.2),
                "Return value:", font_name=FONT_TITLE, font_size=Pt(9),
                font_color=SLATE, bold=True)
    add_textbox(slide, right_left + Pt(10), rv_top + Pt(22),
                Inches(4.0), Inches(0.8),
                '{\n  "primary_label": "retrieval_failure",\n  "evidence_refs": ["span:abc", "artifact:def"]\n}',
                font_name=FONT_MONO, font_size=Pt(9), font_color=INK_NAVY)

    # Stamps
    add_stamp(slide, right_left + Inches(3.2), rv_top + Pt(4), "SCHEMA OK", GREEN)
    add_stamp(slide, right_left + Inches(3.2), rv_top + Pt(24))

    add_footer(slide, 3)
    add_speaker_notes(slide,
        '"Think of it like a program that calls functions and returns structured '
        'values, not a chat bot generating text."')


def build_slide_4(prs):
    """RLM vs Agentic Tool-Use: What's Actually Different?"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_ledger_lines(slide)
    add_ledger_header(slide)
    add_slide_title(slide, "RLM vs Agentic Tool-Use: What's Actually Different?",
                    subtitle="Structured isolation vs linear transcript growth")

    # LEFT COLUMN: Typical tool agent
    col_left = MARGIN
    add_textbox(slide, col_left, Inches(1.7), Inches(5.8), Inches(0.3),
                "Typical tool agent (linear transcript)",
                font_name=FONT_TITLE, font_size=Pt(13), font_color=SLATE, bold=True)

    # Chat transcript card
    chat_card = add_card(slide, col_left, Inches(2.1), Inches(5.8), Inches(3.8),
                         fill_color=WHITE, border_color=BORDER_TAN)
    chat_lines = [
        ("user:", "What caused the failure?", INK_NAVY),
        ("assistant:", "Let me check the traces...", SLATE),
        ("tool:", "list_spans(trace_id=abc) -> 120 spans", TEAL),
        ("assistant:", "I see errors in retrieval. Let me dig deeper...", SLATE),
        ("tool:", "get_span(span_id=xyz) -> {status: ERROR}", TEAL),
        ("\u2718 retry:", "get_span(span_id=xyz) -> timeout", ACCENT_RED),
        ("assistant:", "The retrieval span failed due to...", SLATE),
        ("\u2718 retry:", "list_spans(trace_id=abc) -> stale", ACCENT_RED),
        ("assistant:", "Based on all evidence I've seen...", SLATE),
    ]
    for i, (prefix, text, color) in enumerate(chat_lines):
        y = Inches(2.2) + Inches(i * 0.38)
        add_textbox(slide, col_left + Pt(12), y, Inches(1.0), Inches(0.3),
                    prefix, font_name=FONT_MONO, font_size=Pt(8),
                    font_color=color, bold=True)
        add_textbox(slide, col_left + Inches(1.0), y, Inches(4.5), Inches(0.3),
                    text, font_name=FONT_MONO, font_size=Pt(8),
                    font_color=color)

    # Context drift risk badge
    add_chip(slide, col_left + Inches(3.5), Inches(5.6), "context drift risk",
             bg_color=RGBColor(0xFD, 0xF0, 0xF0), text_color=ACCENT_RED, font_size=Pt(8))
    add_chip(slide, col_left + Pt(12), Inches(6.05), "context length: 18k tokens",
             text_color=SLATE, font_size=Pt(8))

    # RIGHT COLUMN: RLM runtime
    right_left = Inches(7.0)
    add_textbox(slide, right_left, Inches(1.7), Inches(5.8), Inches(0.3),
                "RLM runtime (program-like)",
                font_name=FONT_TITLE, font_size=Pt(13), font_color=TEAL, bold=True)

    # Frame 0
    frame0_top = Inches(2.1)
    add_card(slide, right_left, frame0_top, Inches(5.8), Inches(1.4),
             fill_color=CARD_BG, border_color=TEAL)
    add_textbox(slide, right_left + Pt(10), frame0_top + Pt(6), Inches(3), Inches(0.2),
                "Frame 0: root investigation", font_name=FONT_TITLE,
                font_size=Pt(10), font_color=TEAL, bold=True)
    add_budget_strip(slide, right_left + Pt(10), frame0_top + Pt(24),
                     "depth 0/2 | iter 1/8 | cost $0.003", width=Inches(3.5))
    add_textbox(slide, right_left + Pt(10), frame0_top + Inches(0.7), Inches(5.3), Inches(0.2),
                'tool: list_spans(trace_id=abc) -> 8 hot spans',
                font_name=FONT_MONO, font_size=Pt(8), font_color=INK_NAVY)

    # Frame 1 (nested)
    frame1_top = Inches(2.1) + Inches(1.0)
    add_card(slide, right_left + Inches(0.3), frame1_top + Inches(0.6), Inches(5.2), Inches(1.2),
             fill_color=RGBColor(0xF5, 0xFB, 0xF7), border_color=TEAL)
    add_textbox(slide, right_left + Inches(0.45), frame1_top + Inches(0.68), Inches(3), Inches(0.2),
                "Frame 1: subcall (retrieval hypothesis)", font_name=FONT_TITLE,
                font_size=Pt(9), font_color=TEAL, bold=True)
    add_budget_strip(slide, right_left + Inches(0.45), frame1_top + Inches(0.92),
                     "depth 1/2 | iter 1/4 | cost $0.001", width=Inches(3.2))
    add_textbox(slide, right_left + Inches(0.45), frame1_top + Inches(1.2), Inches(5.0), Inches(0.2),
                'return: { evidence_refs: ["span:xyz", "artifact:ret-001"] }',
                font_name=FONT_MONO, font_size=Pt(8), font_color=INK_NAVY)

    # Frame 2 placeholder (faint)
    add_card(slide, right_left + Inches(0.6), frame1_top + Inches(1.55), Inches(4.9), Inches(0.5),
             fill_color=RGBColor(0xF8, 0xF8, 0xF5), border_color=RGBColor(0xE0, 0xDC, 0xD4))
    add_textbox(slide, right_left + Inches(0.75), frame1_top + Inches(1.6), Inches(3), Inches(0.3),
                "Frame 2: (depth limit reached)",
                font_name=FONT_MONO, font_size=Pt(8), font_color=RGBColor(0xC0, 0xB8, 0xAC))

    # Pool budget
    add_chip(slide, right_left + Pt(10), Inches(5.6), "shared budget: cost <= $0.05",
             bg_color=RGBColor(0xEE, 0xF8, 0xF0), text_color=GREEN, font_size=Pt(8))

    # Bottom strip: Contracts + audit trail
    strip_y = Inches(6.5)
    add_textbox(slide, MARGIN, strip_y - Inches(0.25), Inches(12), Inches(0.2),
                "Contracts + audit trail",
                font_name=FONT_TITLE, font_size=Pt(10), font_color=SLATE, bold=True)
    for i, label in enumerate(["schema validated", "run_record.json", "subcall metadata"]):
        color = GREEN if "schema" in label else MONO_BG
        tc = GREEN if "schema" in label else INK_NAVY
        add_chip(slide, MARGIN + Inches(i * 2.2), strip_y, label,
                 bg_color=color if "schema" in label else MONO_BG,
                 text_color=WHITE if "schema" in label else INK_NAVY,
                 font_size=Pt(8), width=Inches(1.8))

    add_footer(slide, 4)
    add_speaker_notes(slide,
        'Be honest: "An RLM is still an agentic system, but it is a more constrained, '
        'auditable, hierarchical form of agentic execution."\n'
        'The real differentiator is engineering: isolation, budgets, typed interfaces, and replay artifacts.')


def build_slide_5(prs):
    """Capability 1: Root Cause Analysis (RCA)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_ledger_lines(slide)
    add_ledger_header(slide)
    add_slide_title(slide, "Capability 1: Root Cause Analysis (RCA)")

    # Left bullets
    bullets = [
        (0, "RCA goal: assign primary failure mode with evidence pointers + remediation.",
         {"bold": True}),
        (0, ""),
        (0, "Why it matters:"),
        (1, "reduces MTTR (mean time to resolution)"),
        (1, "standardizes incident classification across teams"),
        (1, "makes investigations reproducible and measurable"),
        (0, ""),
        (0, "Output (structured, evidence-linked):"),
        (1, "primary_label  (tool failure, retrieval failure, schema mismatch, ...)"),
        (1, "evidence_refs  (trace_id, span_id, artifact refs)"),
        (1, "summary + remediation + gaps"),
        (0, ""),
        (0, "\u2192 Click-through evidence in the observability UI (span + artifact)",
         {"font_color": TEAL, "bold": True}),
    ]
    add_bullet_slide_text(slide, MARGIN, Inches(1.5), Inches(6.2), Inches(5.5), bullets)

    # Right: RCA output card
    card_left = Inches(7.2)
    card_top = Inches(1.3)
    add_card(slide, card_left, card_top, Inches(5.6), Inches(3.0))

    # RCAReport header
    add_textbox(slide, card_left + Pt(10), card_top + Pt(8), Inches(2.5), Inches(0.25),
                "RCAReport", font_name=FONT_MONO, font_size=Pt(12),
                font_color=INK_NAVY, bold=True)
    add_stamp(slide, card_left + Inches(3.8), card_top + Pt(6))

    # Primary label tag
    tag = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, card_left + Pt(12), card_top + Inches(0.45),
        Inches(1.8), Inches(0.3)
    )
    tag.fill.solid()
    tag.fill.fore_color.rgb = ACCENT_RED
    tag.line.fill.background()
    tf = tag.text_frame
    tf.paragraphs[0].text = "retrieval_failure"
    tf.paragraphs[0].font.name = FONT_MONO
    tf.paragraphs[0].font.size = Pt(11)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.margin_top = Pt(3)

    # Confidence meter
    add_textbox(slide, card_left + Inches(2.2), card_top + Inches(0.48),
                Inches(2), Inches(0.2),
                "confidence: HIGH  \u2588\u2588\u2588\u2591",
                font_name=FONT_MONO, font_size=Pt(9), font_color=GREEN)

    # Evidence chips
    chips_data = [
        ("trace_id: abc-12345678", Inches(0)),
        ("span_id: ret-span-001", Inches(0.32)),
        ("span_id: tool-span-003", Inches(0.64)),
        ("artifact: retrieval-io-7", Inches(0.96)),
    ]
    for label, y_off in chips_data:
        add_chip(slide, card_left + Pt(12), card_top + Inches(0.9) + y_off,
                 label, font_size=Pt(8), width=Inches(2.4))

    # Remediation
    add_textbox(slide, card_left + Inches(2.8), card_top + Inches(0.95),
                Inches(2.6), Inches(0.8),
                "Remediation:\n\u2022 Fix retriever embedding model mismatch\n\u2022 Add citation validation to RAG pipeline",
                font_name=FONT_BODY, font_size=Pt(9), font_color=INK_NAVY)

    # Gaps
    add_textbox(slide, card_left + Inches(2.8), card_top + Inches(2.0),
                Inches(2.6), Inches(0.3),
                "Gaps: upstream latency data not available in trace",
                font_name=FONT_BODY, font_size=Pt(8), font_color=SLATE, italic=True)

    # Bottom: Hot spans ranking table
    hs_top = Inches(4.5)
    add_textbox(slide, card_left + Pt(10), hs_top - Inches(0.25), Inches(3), Inches(0.2),
                "Hot Spans (ranked)", font_name=FONT_TITLE, font_size=Pt(10),
                font_color=INK_NAVY, bold=True)

    hs_rows = [
        ["#", "Span", "Signal", "Why it matters"],
        ["1", "ret-span-001", "ERROR", "retriever returned 0 chunks"],
        ["2", "tool-span-003", "ERROR", "tool input schema mismatch"],
        ["3", "synth-span-005", "SLOW", "synthesis latency > 5s"],
    ]
    add_simple_table(slide, card_left, hs_top, Inches(5.6), Inches(1.5), hs_rows,
                     col_widths=[Inches(0.4), Inches(1.4), Inches(0.8), Inches(3.0)])

    # Trace mini-map (simple tree)
    mm_left = Inches(10.8)
    mm_top = Inches(4.5)
    add_textbox(slide, mm_left, mm_top - Inches(0.25), Inches(2), Inches(0.2),
                "Trace Tree", font_name=FONT_TITLE, font_size=Pt(9),
                font_color=SLATE, bold=True)
    tree_labels = [
        ("root", SLATE, 0),
        ("  \u251c tool-call", SLATE, 0),
        ("  \u251c retrieval", ACCENT_RED, 1),
        ("  \u2502   \u2514 chunk-eval", ACCENT_RED, 1),
        ("  \u2514 synthesis", SLATE, 0),
    ]
    for i, (label, color, highlighted) in enumerate(tree_labels):
        font_color = TEAL if highlighted else color
        add_textbox(slide, mm_left, mm_top + Inches(i * 0.35), Inches(2), Inches(0.25),
                    label, font_name=FONT_MONO, font_size=Pt(8),
                    font_color=font_color, bold=bool(highlighted))

    add_footer(slide, 5)
    add_speaker_notes(slide,
        'Emphasize "evidence linked to trace spans" so engineers can click through in observability UI.')


def build_slide_6(prs):
    """Capability 2: Policy-to-Trace Compliance."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_ledger_lines(slide)
    add_ledger_header(slide)
    add_slide_title(slide, "Capability 2: Policy-to-Trace Compliance")

    # Left bullets
    bullets = [
        (0, "Compliance goal: evaluate behavior against explicit controls using trace evidence.",
         {"bold": True}),
        (0, ""),
        (0, "Why it matters:"),
        (1, "governance: controls need provable evidence, not just narrative"),
        (1, "audit: structured decisions with evidence pointers"),
        (1, "standardization: consistent evaluation across apps/teams"),
        (0, ""),
        (0, "Output (structured, per-control + overall):"),
        (1, "verdict: pass / fail / needs_review / insufficient_evidence"),
        (1, "confidence + severity + missing evidence list"),
        (1, "evidence refs mapped to spans/artifacts"),
    ]
    add_bullet_slide_text(slide, MARGIN, Inches(1.5), Inches(6.2), Inches(5.0), bullets)

    # Right: Overall verdict header
    card_left = Inches(7.2)
    card_top = Inches(1.3)

    # Verdict header card
    vh = add_card(slide, card_left, card_top, Inches(5.6), Inches(0.45),
                  fill_color=RGBColor(0xEE, 0xF8, 0xF0), border_color=GREEN)
    add_textbox(slide, card_left + Pt(12), card_top + Pt(6), Inches(4), Inches(0.3),
                "Overall verdict: PASS     confidence: 0.87",
                font_name=FONT_TITLE, font_size=Pt(13), font_color=GREEN, bold=True)

    # Control ledger table
    tbl_top = card_top + Inches(0.6)
    ctrl_rows = [
        ["control_id", "verdict", "conf", "evidence refs", "\u2713"],
        ["CTRL-001", "PASS", "0.95", "span:tool-001, artifact:io-7", "\u2713"],
        ["CTRL-002", "PASS", "0.88", "span:ret-003, span:synth-005", "\u2713"],
        ["CTRL-003", "NEEDS_REVIEW", "0.62", "span:tool-001", ""],
        ["CTRL-004", "PASS", "0.91", "span:ret-003, artifact:cfg-2", "\u2713"],
        ["CTRL-005", "INSUFF_EVIDENCE", "0.30", "(missing)", ""],
    ]
    add_simple_table(slide, card_left, tbl_top, Inches(5.6), Inches(2.8), ctrl_rows,
                     col_widths=[Inches(1.0), Inches(1.2), Inches(0.5), Inches(2.4), Inches(0.5)])

    # Missing evidence card
    me_left = card_left
    me_top = tbl_top + Inches(3.0)
    add_card(slide, me_left, me_top, Inches(2.8), Inches(1.2),
             fill_color=RGBColor(0xFD, 0xF8, 0xF0), border_color=AMBER)
    add_textbox(slide, me_left + Pt(10), me_top + Pt(6), Inches(2.6), Inches(0.2),
                "Missing evidence", font_name=FONT_TITLE, font_size=Pt(10),
                font_color=AMBER, bold=True)
    missing_items = "\u2610 tool_output artifact\n\u2610 retrieval citations\n\u2610 policy attestation"
    add_textbox(slide, me_left + Pt(10), me_top + Pt(26), Inches(2.6), Inches(0.6),
                missing_items, font_name=FONT_MONO, font_size=Pt(8), font_color=SLATE)
    add_textbox(slide, me_left + Pt(10), me_top + Inches(0.9), Inches(2.6), Inches(0.2),
                "missing \u2192 insufficient_evidence",
                font_name=FONT_MONO, font_size=Pt(7), font_color=AMBER)

    # Policy pack card
    pp_left = card_left + Inches(3.2)
    add_card(slide, pp_left, me_top, Inches(2.4), Inches(1.2),
             fill_color=CARD_BG)
    add_textbox(slide, pp_left + Pt(10), me_top + Pt(6), Inches(2.2), Inches(0.2),
                "Policy pack", font_name=FONT_TITLE, font_size=Pt(10),
                font_color=INK_NAVY, bold=True)
    add_textbox(slide, pp_left + Pt(10), me_top + Pt(26), Inches(2.2), Inches(0.5),
                "controls/v1\nscoped by app + tools",
                font_name=FONT_MONO, font_size=Pt(9), font_color=SLATE)

    add_footer(slide, 6)
    add_speaker_notes(slide,
        'Compliance isn\'t "trust the model." It\'s "prove the model\'s decision with trace evidence."')


def build_slide_7(prs):
    """How RLM Tackles Long Context for RCA + Compliance."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_ledger_lines(slide)
    add_ledger_header(slide)
    add_slide_title(slide, "How RLM Tackles Long Context")

    # Left bullets
    bullets = [
        (0, "Step 1: deterministic narrowing", {"bold": True}),
        (1, "hot spans (errors, exceptions, latency) / scoped controls"),
        (1, "reduces noise and prevents context poisoning"),
        (0, ""),
        (0, "Step 2: evidence-driven recursive loop", {"bold": True}),
        (1, "fetch only needed evidence via tools"),
        (1, "delegate subcalls for subproblems"),
        (0, ""),
        (0, "Step 3: finalize into contracts", {"bold": True}),
        (1, "schema-validated JSON output"),
        (1, "run record emitted for replay and audit"),
    ]
    add_bullet_slide_text(slide, MARGIN, Inches(1.5), Inches(5.8), Inches(5.0), bullets)

    # Right visual: funnel-to-tree
    right_left = Inches(6.8)

    # Raw evidence pile (scattered chips)
    add_textbox(slide, right_left + Inches(0.5), Inches(1.2), Inches(5), Inches(0.2),
                "raw evidence (too big for one prompt)",
                font_name=FONT_MONO, font_size=Pt(8), font_color=SLATE)
    scatter_chips = ["span", "tool_io", "chunk", "config", "control",
                     "span", "msg", "timeout", "schema_err", "rate_limit"]
    for i, label in enumerate(scatter_chips):
        x = right_left + Inches((i % 5) * 1.1) + Inches(0.1)
        y = Inches(1.45) + Inches((i // 5) * 0.32)
        add_chip(slide, x, y, label, font_size=Pt(7),
                 bg_color=RGBColor(0xEE, 0xE8, 0xDE), text_color=SLATE,
                 width=Inches(0.9))

    # Stage 1: Deterministic narrowing
    s1_top = Inches(2.4)
    add_card(slide, right_left, s1_top, Inches(5.8), Inches(1.0),
             fill_color=CARD_BG)
    add_textbox(slide, right_left + Pt(10), s1_top + Pt(6), Inches(3), Inches(0.2),
                "Stage 1: Deterministic narrowing", font_name=FONT_TITLE,
                font_size=Pt(11), font_color=INK_NAVY, bold=True)
    # Funnel result
    add_chip(slide, right_left + Pt(12), s1_top + Inches(0.4), "8 hot spans",
             bg_color=RGBColor(0xE8, 0xF5, 0xF3), text_color=TEAL,
             font_size=Pt(9), width=Inches(1.2))
    add_chip(slide, right_left + Inches(1.5), s1_top + Inches(0.4), "6 scoped controls",
             bg_color=RGBColor(0xE8, 0xF5, 0xF3), text_color=TEAL,
             font_size=Pt(9), width=Inches(1.5))
    add_textbox(slide, right_left + Inches(3.3), s1_top + Inches(0.45), Inches(2.3), Inches(0.2),
                "120 spans \u2192 8 hot", font_name=FONT_MONO, font_size=Pt(9),
                font_color=ACCENT_RED, bold=True)

    # Stage 2: Evidence-driven recursion
    s2_top = Inches(3.6)
    add_card(slide, right_left, s2_top, Inches(5.8), Inches(1.5),
             fill_color=CARD_BG)
    add_textbox(slide, right_left + Pt(10), s2_top + Pt(6), Inches(3), Inches(0.2),
                "Stage 2: Evidence-driven recursion", font_name=FONT_TITLE,
                font_size=Pt(11), font_color=INK_NAVY, bold=True)
    # Recursion tree
    tree_nodes = [
        (0, "objective", INK_NAVY),
        (1, "\u251c tool_call", TEAL),
        (1, "\u251c delegate_subcall", TEAL),
        (2, "\u2514 tool_call (focused)", SLATE),
        (1, "\u2514 submit", GREEN),
    ]
    for i, (indent, label, color) in enumerate(tree_nodes):
        x = right_left + Pt(16) + Inches(indent * 0.25)
        y = s2_top + Inches(0.35) + Inches(i * 0.22)
        add_textbox(slide, x, y, Inches(3), Inches(0.2),
                    label, font_name=FONT_MONO, font_size=Pt(9),
                    font_color=color, bold=(i == 0))
    add_budget_strip(slide, right_left + Inches(3.2), s2_top + Inches(0.4),
                     "depth <= 2 | iter <= 8", width=Inches(2.3))

    # Stage 3: Typed finalize
    s3_top = Inches(5.3)
    add_card(slide, right_left, s3_top, Inches(5.8), Inches(0.8),
             fill_color=CARD_BG)
    add_textbox(slide, right_left + Pt(10), s3_top + Pt(6), Inches(3), Inches(0.2),
                "Stage 3: Typed finalize", font_name=FONT_TITLE,
                font_size=Pt(11), font_color=INK_NAVY, bold=True)
    add_stamp(slide, right_left + Pt(12), s3_top + Inches(0.35), "SCHEMA OK", GREEN)
    add_stamp(slide, right_left + Inches(1.6), s3_top + Inches(0.35))

    # Evidence chips attached to finalize
    for i, label in enumerate(["trace_id: abc", "span_id: ret-001", "artifact: io-7"]):
        add_chip(slide, right_left + Inches(3.0) + Inches(i * 1.0), s3_top + Inches(0.38),
                 label, font_size=Pt(7), width=Inches(0.95))

    # Bounded cost badge
    add_chip(slide, right_left + Inches(0.5), Inches(6.3), "bounded cost/run: cost <= $0.05",
             bg_color=RGBColor(0xEE, 0xF8, 0xF0), text_color=GREEN, font_size=Pt(9),
             width=Inches(2.8))

    add_footer(slide, 7)
    add_speaker_notes(slide,
        'This is the core "long context automation" idea: search the evidence space '
        'instead of stuffing tokens.')


def build_slide_8(prs):
    """Implementation: What The Investigation Layer Is Doing Under the Hood."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_ledger_lines(slide)
    add_ledger_header(slide)
    add_slide_title(slide, "Implementation: Under the Hood")

    # Left bullets
    bullets = [
        (0, "Investigator runtime is a bounded REPL loop:", {"bold": True}),
        (1, "call_tool(...)  = read-only inspection API fetch"),
        (1, "llm_query(...)  = short semantic synthesis calls"),
        (1, "SUBMIT(...)     = finalize structured output"),
        (0, ""),
        (0, "Safety:"),
        (1, "sandbox tool allowlist + argument validation"),
        (1, "budgets (iterations, depth, tools, tokens, cost, wall-time)"),
        (1, "deterministic recovery submit near deadline"),
        (0, ""),
        (0, "Auditability:"),
        (1, "every run writes: run_record.json + output JSON"),
        (1, "trajectory, usage/cost, evidence pointers, errors"),
    ]
    add_bullet_slide_text(slide, MARGIN, Inches(1.5), Inches(5.8), Inches(5.0), bullets)

    # Right: Architecture row
    right_left = Inches(6.8)
    arch_top = Inches(1.2)

    # Four architecture cards with arrows
    arch_items = [
        ("Trace store\n+ artifacts", Inches(0)),
        ("Investigation\nqueue", Inches(1.55)),
        ("Investigator\nworker", Inches(3.1)),
        ("Write-back\nfindings", Inches(4.65)),
    ]
    for label, x_off in arch_items:
        add_card(slide, right_left + x_off, arch_top, Inches(1.3), Inches(0.7),
                 fill_color=CARD_BG)
        add_textbox(slide, right_left + x_off + Pt(6), arch_top + Pt(6),
                    Inches(1.2), Inches(0.55),
                    label, font_name=FONT_BODY, font_size=Pt(8),
                    font_color=INK_NAVY, alignment=PP_ALIGN.CENTER)

    # Arrows between architecture cards
    for i in range(3):
        x1 = right_left + arch_items[i][1] + Inches(1.3)
        x2 = right_left + arch_items[i + 1][1]
        y = arch_top + Inches(0.35)
        add_textbox(slide, x1, y - Pt(4), x2 - x1, Inches(0.25),
                    "\u2192", font_name=FONT_MONO, font_size=Pt(14),
                    font_color=BORDER_TAN, alignment=PP_ALIGN.CENTER)

    # Runtime internals card
    rt_top = Inches(2.15)
    add_card(slide, right_left, rt_top, Inches(5.8), Inches(1.1),
             fill_color=RGBColor(0xF8, 0xF6, 0xF1), border_color=TEAL)
    add_textbox(slide, right_left + Pt(10), rt_top + Pt(4), Inches(3), Inches(0.2),
                "RLM REPL Runtime", font_name=FONT_TITLE, font_size=Pt(12),
                font_color=TEAL, bold=True)
    # Sub-badges
    badges = [
        ("ToolRegistry", "allowlisted"),
        ("Sandbox", "arg checks"),
        ("Budget", "depth/iter/cost/wall"),
        ("RunRecord", "artifacts"),
    ]
    for i, (name, desc) in enumerate(badges):
        bx = right_left + Inches(i * 1.4) + Pt(10)
        by = rt_top + Inches(0.4)
        add_chip(slide, bx, by, name, bg_color=MONO_BG, text_color=INK_NAVY,
                 font_size=Pt(8), width=Inches(1.2))
        add_textbox(slide, bx, by + Inches(0.3), Inches(1.2), Inches(0.2),
                    desc, font_name=FONT_MONO, font_size=Pt(7),
                    font_color=SLATE, alignment=PP_ALIGN.CENTER)

    # REPL panel (code cell)
    repl_top = Inches(3.5)
    add_card(slide, right_left, repl_top, Inches(3.5), Inches(2.2),
             fill_color=RGBColor(0x1B, 0x2A, 0x3A), border_color=TEAL)
    repl_lines = [
        '# REPL iteration 3/8',
        'spans = call_tool("list_spans",',
        '    trace_id="abc-123")',
        'hot = [s for s in spans',
        '    if s.status == "ERROR"]',
        'evidence = call_tool("get_tool_io",',
        '    span_id=hot[0].id)',
        'SUBMIT({"primary_label": ...,',
        '    "evidence_refs": [...]})',
    ]
    for i, line in enumerate(repl_lines):
        add_textbox(slide, right_left + Pt(10), repl_top + Pt(8) + Inches(i * 0.22),
                    Inches(3.3), Inches(0.2),
                    line, font_name=FONT_MONO, font_size=Pt(8),
                    font_color=RGBColor(0xA8, 0xDB, 0xAF))
    add_chip(slide, right_left + Inches(2.5), repl_top + Inches(1.9),
             "wall-time 18s", bg_color=RGBColor(0x2A, 0x3A, 0x4A),
             text_color=RGBColor(0xA8, 0xDB, 0xAF), font_size=Pt(7))

    # Tool allowlist strip
    tools_top = Inches(3.5)
    tools_left = right_left + Inches(3.8)
    add_textbox(slide, tools_left, tools_top, Inches(2), Inches(0.2),
                "Tool allowlist", font_name=FONT_TITLE, font_size=Pt(9),
                font_color=INK_NAVY, bold=True)
    tool_items = [
        ("list_spans", "trace graph"),
        ("get_span", "details"),
        ("get_tool_io", "tool evidence"),
        ("list_controls", "policy pack"),
        ("required_evidence", "control needs"),
        ("get_retrieval_io", "RAG context"),
    ]
    for i, (name, desc) in enumerate(tool_items):
        add_chip(slide, tools_left, tools_top + Inches(0.3) + Inches(i * 0.32),
                 name, font_size=Pt(7), width=Inches(1.6))
        add_textbox(slide, tools_left + Inches(0.1), tools_top + Inches(0.58) + Inches(i * 0.32),
                    Inches(1.6), Inches(0.15),
                    desc, font_name=FONT_MONO, font_size=Pt(6),
                    font_color=SLATE, alignment=PP_ALIGN.CENTER)

    # Auditability footer
    af_y = Inches(6.3)
    for i, label in enumerate(["state_trajectory", "usage/cost", "subcall_metadata"]):
        add_chip(slide, right_left + Inches(i * 2.0), af_y, label,
                 font_size=Pt(8), width=Inches(1.7))

    add_footer(slide, 8)
    add_speaker_notes(slide,
        '"This is built for operations: you can inspect the execution trajectory and costs for each run."')


def build_slide_9(prs):
    """Results: Proof Canaries."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_ledger_lines(slide)
    add_ledger_header(slide)
    add_slide_title(slide, "Results: Proof Canaries (Quality + Run Health + Cost)")

    # Left bullets
    bullets = [
        (0, "RCA 5-trace canary:", {"bold": True}),
        (1, "baseline accuracy 0.0 \u2192 RLM 0.8  (delta +0.8)"),
        (1, "run health: succeeded 5/5, wall-time partials 0"),
        (1, "cost: total ~$0.0216  (avg ~$0.0043/trace)"),
        (0, ""),
        (0, "Compliance 5-trace canary:", {"bold": True}),
        (1, "baseline accuracy 0.8 \u2192 RLM 1.0  (delta +0.2)"),
        (1, "run health: succeeded 5/5, partial_rate 0.0, failed 0"),
        (1, "cost: total ~$0.0469  (avg ~$0.0094/trace)"),
        (0, ""),
        (0, "Key takeaway: working beta with strong run-health\nand measurable quality lift.",
         {"bold": True, "font_color": TEAL}),
    ]
    add_bullet_slide_text(slide, MARGIN, Inches(1.5), Inches(5.8), Inches(5.0), bullets)

    # Right: Scorecard table
    card_left = Inches(6.8)
    card_top = Inches(1.3)

    score_rows = [
        ["Capability", "Baseline", "RLM", "Delta", "Health", "Avg Cost"],
        ["RCA", "0.0", "0.8", "+0.8", "S5 P0 F0", "$0.0043"],
        ["Compliance", "0.8", "1.0", "+0.2", "S5 P0 F0", "$0.0094"],
    ]
    tbl = add_simple_table(slide, card_left, card_top, Inches(5.8), Inches(1.3), score_rows,
                           col_widths=[Inches(1.2), Inches(0.8), Inches(0.7), Inches(0.7),
                                       Inches(1.1), Inches(1.0)])

    # Delta visual bars
    for i, (delta_val, delta_width) in enumerate([(0.8, 3.2), (0.2, 0.8)]):
        bar_y = card_top + Inches(0.6) + Inches(i * 0.45)
        bar_x = card_left + Inches(0.3)
        # Delta bar background
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, card_left, bar_y + Inches(1.2),
            Inches(delta_width), Inches(0.15)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = TEAL
        bar.line.fill.background()

    # Delta labels
    add_textbox(slide, card_left, card_top + Inches(1.85), Inches(2), Inches(0.25),
                "RCA delta: +0.8", font_name=FONT_MONO, font_size=Pt(14),
                font_color=TEAL, bold=True)
    add_textbox(slide, card_left + Inches(3), card_top + Inches(1.85), Inches(2), Inches(0.25),
                "Compliance delta: +0.2", font_name=FONT_MONO, font_size=Pt(14),
                font_color=TEAL, bold=True)

    # Run health tiles
    rh_top = Inches(3.7)
    add_textbox(slide, card_left, rh_top, Inches(3), Inches(0.2),
                "Run Health Summary", font_name=FONT_TITLE, font_size=Pt(12),
                font_color=INK_NAVY, bold=True)

    for i, (cap, s, p, f) in enumerate([("RCA", 5, 0, 0), ("Compliance", 5, 0, 0)]):
        y = rh_top + Inches(0.35) + Inches(i * 0.45)
        add_textbox(slide, card_left, y, Inches(1.2), Inches(0.3),
                    cap, font_name=FONT_TITLE, font_size=Pt(10),
                    font_color=INK_NAVY, bold=True)
        add_chip(slide, card_left + Inches(1.3), y, f"S {s}", width=Inches(0.6),
                 bg_color=RGBColor(0xEE, 0xF8, 0xF0), text_color=GREEN, font_size=Pt(10))
        add_chip(slide, card_left + Inches(2.0), y, f"P {p}", width=Inches(0.6),
                 bg_color=MONO_BG, text_color=SLATE, font_size=Pt(10))
        add_chip(slide, card_left + Inches(2.7), y, f"F {f}", width=Inches(0.6),
                 bg_color=MONO_BG, text_color=SLATE, font_size=Pt(10))

    # Combined cost
    add_textbox(slide, card_left, Inches(4.9), Inches(5), Inches(0.3),
                "Combined avg: ~$0.0137/trace  (RCA + Compliance)",
                font_name=FONT_MONO, font_size=Pt(12), font_color=INK_NAVY, bold=True)

    # Evidence callout
    ev_top = Inches(5.5)
    add_textbox(slide, card_left, ev_top, Inches(2), Inches(0.2),
                "Evidence", font_name=FONT_TITLE, font_size=Pt(10),
                font_color=INK_NAVY, bold=True)
    add_chip(slide, card_left, ev_top + Inches(0.28),
             "rca_only_report.json", font_size=Pt(7), width=Inches(2.2))
    add_chip(slide, card_left, ev_top + Inches(0.62),
             "compliance_only_report.json", font_size=Pt(7), width=Inches(2.5))
    add_stamp(slide, card_left + Inches(2.8), ev_top + Inches(0.35))

    add_footer(slide, 9)
    add_speaker_notes(slide,
        'If asked "why baseline is low": baseline is a simple deterministic heuristic; '
        'we use it only to measure non-regression and delta.')


def build_slide_10(prs):
    """Scaling to Company-Wide Observability."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_ledger_lines(slide)
    add_ledger_header(slide)
    add_slide_title(slide, "Scaling to Company-Wide Observability")

    # Left bullets
    bullets = [
        (0, "Deployment model:", {"bold": True}),
        (1, "async job queue triggered by traces/incidents"),
        (1, "investigator workers per capability (RCA, compliance)"),
        (1, "write-back findings to observability platform UI"),
        (0, ""),
        (0, "Guarded rollout:", {"bold": True}),
        (1, "policy packs per domain/team"),
        (1, "budgets per run + per team + per day"),
        (1, "route low-confidence/high-severity to humans"),
        (1, "continuous benchmarking gates (quality + run-health + cost)"),
        (0, ""),
        (0, "Economics (pilot profile):", {"bold": True}),
        (1, "combined avg ~$0.0137/trace (RCA + compliance)"),
        (1, "scale by automating repetitive, high-volume patterns"),
    ]
    add_bullet_slide_text(slide, MARGIN, Inches(1.5), Inches(5.8), Inches(5.0), bullets)

    # Right: Two-lane operating model
    right_left = Inches(6.8)

    # Lane A: Automation
    la_top = Inches(1.3)
    add_card(slide, right_left, la_top, Inches(5.8), Inches(1.5),
             fill_color=RGBColor(0xF2, 0xF9, 0xF3), border_color=GREEN)
    add_textbox(slide, right_left + Pt(10), la_top + Pt(4), Inches(3), Inches(0.2),
                "Lane A: Automation (bounded)", font_name=FONT_TITLE,
                font_size=Pt(11), font_color=GREEN, bold=True)
    # Flow chips
    flow_items = ["Trace\ntrigger", "Queue", "Worker\npool", "Findings", "Write-\nback"]
    for i, label in enumerate(flow_items):
        x = right_left + Inches(i * 1.15) + Pt(10)
        add_chip(slide, x, la_top + Inches(0.4), label.replace("\n", " "),
                 bg_color=WHITE, text_color=GREEN, font_size=Pt(7), width=Inches(0.9))
        if i < len(flow_items) - 1:
            add_textbox(slide, x + Inches(0.92), la_top + Inches(0.4), Inches(0.2), Inches(0.28),
                        "\u2192", font_name=FONT_MONO, font_size=Pt(10), font_color=GREEN)
    # Budget envelope
    add_chip(slide, right_left + Pt(10), la_top + Inches(0.85),
             "daily cap | per-run budget | wall-time",
             bg_color=MONO_BG, text_color=SLATE, font_size=Pt(7), width=Inches(3.0))
    # Policy packs + benchmark gates
    add_chip(slide, right_left + Inches(3.3), la_top + Inches(0.85),
             "policy packs \u2192 worker",
             bg_color=MONO_BG, text_color=INK_NAVY, font_size=Pt(7), width=Inches(1.7))

    # Lane B: Human review
    lb_top = Inches(3.0)
    add_card(slide, right_left, lb_top, Inches(5.8), Inches(1.1),
             fill_color=RGBColor(0xFD, 0xF8, 0xF0), border_color=AMBER)
    add_textbox(slide, right_left + Pt(10), lb_top + Pt(4), Inches(4), Inches(0.2),
                "Lane B: Human review (only when needed)", font_name=FONT_TITLE,
                font_size=Pt(11), font_color=AMBER, bold=True)
    add_textbox(slide, right_left + Pt(10), lb_top + Inches(0.35), Inches(5.5), Inches(0.3),
                "low confidence / high severity  \u2192  analyst  \u2192  approval / escalation",
                font_name=FONT_MONO, font_size=Pt(9), font_color=SLATE)
    add_chip(slide, right_left + Pt(10), lb_top + Inches(0.7),
             "evidence pack for analyst",
             bg_color=MONO_BG, text_color=INK_NAVY, font_size=Pt(7), width=Inches(2.2))

    # Benchmark gates
    add_chip(slide, right_left + Inches(3.5), lb_top + Inches(0.7),
             "benchmark gates: quality + run health + cost",
             bg_color=RGBColor(0xEE, 0xF8, 0xF0), text_color=GREEN, font_size=Pt(7),
             width=Inches(3.0))

    # Rollout timeline
    rt_top = Inches(4.5)
    add_textbox(slide, right_left, rt_top, Inches(3), Inches(0.2),
                "Rollout Timeline", font_name=FONT_TITLE, font_size=Pt(12),
                font_color=INK_NAVY, bold=True)

    milestones = [
        ("1. Canary", "N small", "non-negative delta"),
        ("2. Guarded rollout", "team opt-in", "budget caps"),
        ("3. Standard workflow", "org-wide", "SLA/MTTR impact"),
    ]
    for i, (title, scope, metric) in enumerate(milestones):
        x = right_left + Inches(i * 2.0)
        y = rt_top + Inches(0.35)
        add_card(slide, x, y, Inches(1.8), Inches(1.1), fill_color=CARD_BG)
        add_textbox(slide, x + Pt(8), y + Pt(6), Inches(1.6), Inches(0.2),
                    title, font_name=FONT_TITLE, font_size=Pt(10),
                    font_color=INK_NAVY, bold=True)
        add_textbox(slide, x + Pt(8), y + Pt(26), Inches(1.6), Inches(0.2),
                    scope, font_name=FONT_MONO, font_size=Pt(8), font_color=SLATE)
        add_textbox(slide, x + Pt(8), y + Pt(44), Inches(1.6), Inches(0.2),
                    metric, font_name=FONT_MONO, font_size=Pt(8), font_color=TEAL, bold=True)

    # Arrows between milestones
    for i in range(2):
        x = right_left + Inches((i + 1) * 2.0) - Inches(0.15)
        add_textbox(slide, x, rt_top + Inches(0.7), Inches(0.3), Inches(0.3),
                    "\u2192", font_name=FONT_MONO, font_size=Pt(14),
                    font_color=BORDER_TAN, alignment=PP_ALIGN.CENTER)

    add_footer(slide, 10)
    add_speaker_notes(slide,
        'Emphasize "this becomes an observability investigation layer" not just a one-off agent.')


def build_slide_11(prs):
    """Roadmap: What Comes Next."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_ledger_lines(slide)
    add_ledger_header(slide)
    add_slide_title(slide, "Roadmap: What Comes Next (Credible, Not Hype)")

    # Left bullets
    bullets = [
        (0, "Improve RCA edge cases"),
        (1, "e.g., tool timeout vs upstream dependency classification"),
        (0, "Enforce stricter RCA REPL submit contract"),
        (1, "prevent missing primary_label"),
        (0, "Add connectors for logs/metrics (beyond traces)"),
        (0, "Expand incident dossier engine"),
        (1, "once multi-trace correlation tools are ready"),
    ]
    add_bullet_slide_text(slide, MARGIN, Inches(1.5), Inches(5.8), Inches(4.0), bullets)

    # Right: Now / Next / Later columns
    right_left = Inches(6.8)
    col_w = Inches(1.85)

    columns = [
        ("NOW", INK_NAVY, [
            ("disambiguate tool\ntimeout vs upstream", "accuracy gap",
             "delta >= 0 on canaries"),
            ("enforce submit\ncontract", "prevent missing fields",
             "0 missing primary_label"),
        ]),
        ("NEXT", TEAL, [
            ("add logs/metrics\nconnectors", "richer evidence set",
             "new signal types in runs"),
            ("multi-trace\ncorrelation", "incident scope",
             "cross-trace coverage"),
        ]),
        ("LATER", SLATE, [
            ("expand incident\ndossier engine", "full incident product",
             "dossier schema v2"),
            ("self-improving\nproof canaries", "continuous quality",
             "automated regression"),
        ]),
    ]

    for col_i, (header, color, items) in enumerate(columns):
        x = right_left + Inches(col_i * 2.0)
        # Column header
        header_card = add_card(slide, x, Inches(1.3), col_w, Inches(0.4),
                               fill_color=color, border_color=color)
        tf = header_card.text_frame
        tf.paragraphs[0].text = header
        tf.paragraphs[0].font.name = FONT_TITLE
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.margin_top = Pt(4)

        # Item cards
        for item_i, (title, why, metric) in enumerate(items):
            iy = Inches(1.85) + Inches(item_i * 1.8)
            item_card = add_card(slide, x, iy, col_w, Inches(1.55), fill_color=CARD_BG)

            add_textbox(slide, x + Pt(8), iy + Pt(6), col_w - Pt(16), Inches(0.5),
                        title, font_name=FONT_TITLE, font_size=Pt(10),
                        font_color=INK_NAVY, bold=True)
            add_textbox(slide, x + Pt(8), iy + Inches(0.55), col_w - Pt(16), Inches(0.2),
                        why, font_name=FONT_BODY, font_size=Pt(8), font_color=SLATE)
            add_textbox(slide, x + Pt(8), iy + Inches(0.8), col_w - Pt(16), Inches(0.2),
                        metric, font_name=FONT_MONO, font_size=Pt(8), font_color=TEAL)
            add_chip(slide, x + Pt(8), iy + Inches(1.1), "proof_report.json",
                     font_size=Pt(6), width=Inches(1.3))

    # Bottom rule banner
    banner_y = Inches(6.3)
    banner = add_card(slide, MARGIN, banner_y, Inches(12.3), Inches(0.4),
                      fill_color=MONO_BG, border_color=BORDER_TAN)
    add_textbox(slide, MARGIN + Pt(20), banner_y + Pt(4), Inches(11), Inches(0.3),
                "Rule: improvements must show non-negative delta on proof canaries",
                font_name=FONT_MONO, font_size=Pt(10), font_color=ACCENT_RED, bold=True,
                alignment=PP_ALIGN.CENTER)

    add_footer(slide, 11)
    add_speaker_notes(slide,
        'Make clear the difference between "working runtime" and "perfect accuracy." '
        'We have a measurable path to improve.')


def build_slide_12(prs):
    """References (External + Internal)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_ledger_lines(slide)
    add_ledger_header(slide)
    add_slide_title(slide, "References (External + Internal)")

    # Left: External references as citation cards
    refs = [
        ("Reasoning Language Models (RLM)", "Core architecture pattern for bounded investigation",
         "arxiv", "arxiv.org/abs/2512.24601"),
        ("Lost in the Middle", "Why long context degrades accuracy for retrieval tasks",
         "arxiv", "arxiv.org/abs/2307.03172"),
        ("OpenTelemetry Signals", "Trace/span/metric data model for observability",
         "spec", "opentelemetry.io/docs/concepts/signals"),
        ("NIST AI Risk Management Framework", "Enterprise AI governance baseline",
         "spec", "nist.gov/itl/ai-risk-management-framework"),
        ("Arize Phoenix", "Open-source trace store and observability UI",
         "tool", "arize.com/docs/phoenix"),
        ("RLM Explainer (Alex Zhang)", "Practical walkthrough of RLM concepts",
         "concept", "alexzhang13.github.io/blog/2025/rlm"),
    ]

    for i, (title, why, tag, url) in enumerate(refs):
        y = Inches(1.5) + Inches(i * 0.88)
        add_card(slide, MARGIN, y, Inches(6.0), Inches(0.78), fill_color=CARD_BG)
        add_textbox(slide, MARGIN + Pt(10), y + Pt(4), Inches(5.0), Inches(0.2),
                    title, font_name=FONT_TITLE, font_size=Pt(10),
                    font_color=INK_NAVY, bold=True)
        add_textbox(slide, MARGIN + Pt(10), y + Pt(22), Inches(5.0), Inches(0.2),
                    why, font_name=FONT_BODY, font_size=Pt(8), font_color=SLATE)
        # Tag chip
        tag_colors = {"arxiv": ACCENT_RED, "spec": TEAL, "tool": GREEN, "concept": AMBER}
        add_chip(slide, MARGIN + Inches(4.5), y + Pt(4), tag,
                 bg_color=tag_colors.get(tag, MONO_BG),
                 text_color=WHITE, font_size=Pt(7), width=Inches(0.7))
        # URL in small text
        add_textbox(slide, MARGIN + Pt(10), y + Pt(40), Inches(5.5), Inches(0.2),
                    url, font_name=FONT_MONO, font_size=Pt(7), font_color=SLATE)

    # Right: Internal proof index
    right_left = Inches(7.0)
    pi_top = Inches(1.5)
    add_card(slide, right_left, pi_top, Inches(5.5), Inches(3.0), fill_color=CARD_BG)
    add_textbox(slide, right_left + Pt(12), pi_top + Pt(8), Inches(3), Inches(0.3),
                "Proof Index", font_name=FONT_TITLE, font_size=Pt(14),
                font_color=INK_NAVY, bold=True)
    add_stamp(slide, right_left + Inches(3.5), pi_top + Pt(8))

    add_textbox(slide, right_left + Pt(12), pi_top + Inches(0.6), Inches(5), Inches(0.2),
                "RCA canary report:", font_name=FONT_BODY, font_size=Pt(9),
                font_color=SLATE, bold=True)
    add_chip(slide, right_left + Pt(12), pi_top + Inches(0.85),
             "phase10-rca-only-canary-5-retrieverfix/.../rca_only_report.json",
             font_size=Pt(7), width=Inches(5.0))

    add_textbox(slide, right_left + Pt(12), pi_top + Inches(1.3), Inches(5), Inches(0.2),
                "Compliance canary report:", font_name=FONT_BODY, font_size=Pt(9),
                font_color=SLATE, bold=True)
    add_chip(slide, right_left + Pt(12), pi_top + Inches(1.55),
             "phase10-compliance-only-canary-5/.../compliance_only_report.json",
             font_size=Pt(7), width=Inches(5.0))

    add_textbox(slide, right_left + Pt(12), pi_top + Inches(2.1), Inches(5), Inches(0.3),
                "replayable from artifacts",
                font_name=FONT_MONO, font_size=Pt(10), font_color=TEAL, bold=True)
    add_stamp(slide, right_left + Inches(2.5), pi_top + Inches(2.1))

    # Thank you / contact area
    add_textbox(slide, right_left, Inches(5.2), Inches(5.5), Inches(0.5),
                "Thank you",
                font_name=FONT_TITLE, font_size=Pt(24), font_color=INK_NAVY,
                bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, right_left, Inches(5.7), Inches(5.5), Inches(0.3),
                "Questions? Evidence is in the artifacts.",
                font_name=FONT_BODY, font_size=Pt(12), font_color=SLATE,
                alignment=PP_ALIGN.CENTER)

    add_footer(slide, 12)
    add_speaker_notes(slide,
        "External references are in the speaker notes / appendix. "
        "Internal proof reports are replayable from the artifacts directory.")


# ===========================================================================
# Main
# ===========================================================================

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    build_slide_1(prs)
    build_slide_2(prs)
    build_slide_3(prs)
    build_slide_4(prs)
    build_slide_5(prs)
    build_slide_6(prs)
    build_slide_7(prs)
    build_slide_8(prs)
    build_slide_9(prs)
    build_slide_10(prs)
    build_slide_11(prs)
    build_slide_12(prs)

    out_dir = os.path.dirname(os.path.abspath(__file__))
    out_path = os.path.join(out_dir, "ppevigil_rlm_pepsico_deck_v2.pptx")
    prs.save(out_path)
    print(f"Saved: {out_path}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
